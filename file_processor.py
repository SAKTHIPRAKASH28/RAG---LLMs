import PyPDF2
from collections import Counter
import io
from pptx import Presentation
from docx import Document
from PIL import Image
import numpy as np
from transformers import CLIPProcessor, CLIPModel,WhisperProcessor, WhisperForConditionalGeneration
import librosa

def read_audio(file_bytes):

    processor = WhisperProcessor.from_pretrained("openai/whisper-base")
    model = WhisperForConditionalGeneration.from_pretrained("openai/whisper-base")

    audio, sr = librosa.load(io.BytesIO(file_bytes), sr=16000, mono=True)

    input_features = processor(audio, sampling_rate=16000, return_tensors="pt").input_features

    predicted_ids = model.generate(input_features)

    transcription = processor.batch_decode(predicted_ids, skip_special_tokens=True)

    return f"Audio Transcription:\n{transcription[0]}"

def read_image(file_bytes):
    image = Image.open(io.BytesIO(file_bytes))
    
    # CLIP for general image understanding
    model = CLIPModel.from_pretrained("openai/clip-vit-base-patch32")
    processor = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch32")
    
    inputs = processor(images=image, return_tensors="pt", padding=True, truncation=True)
    image_features = model.get_image_features(**inputs)
    
    # Generate multiple descriptions using CLIP
    possible_descriptions = [
        "This is a photo of",
        "This image contains",
        "This picture shows",
        "The main subject of this image is",
        "This image depicts",
        "The scene in this image is",
        "This photograph captures",
        "The primary focus of this image is",
        "This visual representation includes",
        "The key elements in this image are"
    ]
    
    text_inputs = processor(
        text=possible_descriptions,
        return_tensors="pt",
        padding=True,
        truncation=True
    )
    text_features = model.get_text_features(**text_inputs)
    
    similarities = (image_features @ text_features.T).softmax(dim=-1)
    
    # Get top 3 descriptions
    top_matches = similarities.topk(3)
    descriptions = []
    for i in range(3):
        idx = top_matches.indices[0][i].item()
        score = top_matches.values[0][i].item()
        description = possible_descriptions[idx]
        descriptions.append(f"{description} (confidence: {score:.2f})")
    
    return "Image Analysis:\n" + "\n".join(descriptions)


def read_pdf(file_bytes):
    pdf_stream = io.BytesIO(file_bytes)
    reader = PyPDF2.PdfReader(pdf_stream)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def read_pptx(file_bytes):
    pptx_stream = io.BytesIO(file_bytes)
    prs = Presentation(pptx_stream)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + "\n"
    return text

def read_txt(file_bytes):
    return file_bytes.decode('utf-8')

def read_docx(file_bytes):
    docx_stream = io.BytesIO(file_bytes)
    doc = Document(docx_stream)
    return "\n".join([paragraph.text for paragraph in doc.paragraphs])

def preprocess_text(text):
    chunks = text.split('\n')
    return [chunk.strip() for chunk in chunks if chunk.strip()]

def create_index(chunks):
    word_freq = Counter()
    for chunk in chunks:
        word_freq.update(chunk.lower().split())
    return word_freq, chunks

def retrieve_context(query, index, chunks, k=3):
    query_words = set(query.lower().split())
    chunk_scores = [(i, len(set(chunk.lower().split()) & query_words)) for i, chunk in enumerate(chunks)]
    top_chunks = sorted(chunk_scores, key=lambda x: x[1], reverse=True)[:k]
    return [chunks[i] for i, _ in top_chunks]
